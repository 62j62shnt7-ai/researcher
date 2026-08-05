"""Document parsing and chunking.

Each parser returns a list of blocks: {"text": str, "location": str}.
Blocks are then merged/split into overlapping chunks that keep their location label.
"""
import re
from pathlib import Path

SUPPORTED_EXTS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xlsm", ".csv", ".tsv",
    ".pptx", ".txt", ".md", ".markdown", ".rst", ".html", ".htm", ".log",
}


def detect_format(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    return ext.lstrip(".") if ext in SUPPORTED_EXTS else "unknown"


# ---------------- parsers ----------------

# A page with less text than this is considered scanned/non-searchable -> OCR candidate
MIN_PAGE_CHARS = 40


_rapidocr_engine = None


def _ocr_image(img) -> str:
    """Run RapidOCR on a numpy image, handling both old and new package APIs."""
    res = _rapidocr_engine(img)
    if isinstance(res, tuple):          # rapidocr_onnxruntime: (result, elapse)
        res = res[0]
    txts = getattr(res, "txts", None)   # rapidocr >= 2.x: RapidOCROutput object
    if txts is not None:
        return "\n".join(txts)
    if res:                             # old API: [[box, text, score], ...]
        return "\n".join(item[1] for item in res)
    return ""


def _get_ocr():
    """Returns an ocr(path, page_no) -> str function, or None if no OCR backend exists.

    Default backend (pip-only, installed via requirements-ocr.txt):
      pypdfium2 renders the page to an image, RapidOCR recognizes the text.
      Works with both the 'rapidocr' (new) and 'rapidocr-onnxruntime' (classic) packages.
    Fallback backend (optional): pytesseract + pdf2image (needs Tesseract/Poppler)."""
    try:
        import numpy as np
        import pypdfium2 as pdfium
        try:
            from rapidocr import RapidOCR            # new package name
        except ImportError:
            from rapidocr_onnxruntime import RapidOCR  # classic package name
        global _rapidocr_engine
        if _rapidocr_engine is None:
            _rapidocr_engine = RapidOCR()

        def ocr_page(path: str, page_no: int) -> str:
            try:
                pdf = pdfium.PdfDocument(path)
                try:
                    page = pdf[page_no - 1]
                    img = np.asarray(page.render(scale=200 / 72).to_pil().convert("RGB"))
                finally:
                    pdf.close()
                return _ocr_image(img)
            except Exception:
                return ""

        return ocr_page
    except Exception:
        # covers ImportError AND engine-init failures (e.g. rapidocr installed
        # without a working onnxruntime on too-new Python versions)
        pass

    try:
        import pytesseract
        from pdf2image import convert_from_path
    except Exception:
        return None

    def ocr_page_tess(path: str, page_no: int) -> str:
        try:
            images = convert_from_path(path, dpi=200, first_page=page_no, last_page=page_no)
            return pytesseract.image_to_string(images[0]) if images else ""
        except Exception:
            return ""

    return ocr_page_tess


def _pdf_page_texts(path: str) -> list[str]:
    """Extract text per page; pdfplumber first (better text + tables), pypdf fallback."""
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                try:
                    for tbl in page.extract_tables() or []:
                        rows = [" | ".join((cell or "").strip() for cell in row) for row in tbl if row]
                        t = "\n".join(r for r in rows if r.strip(" |"))
                        if t and t not in text:
                            text += "\n[Table]\n" + t
                except Exception:
                    pass
                texts.append(text)
        if any(t.strip() for t in texts):
            return texts
    except Exception:
        pass

    from pypdf import PdfReader
    reader = PdfReader(path)
    return [(page.extract_text() or "") for page in reader.pages]


def _parse_pdf(path: str) -> tuple[list[dict], int]:
    """Per-page searchability detection: pages with extractable text are used directly;
    sparse/empty pages (scanned) are OCR'd individually when OCR is available."""
    page_texts = _pdf_page_texts(path)
    n_pages = len(page_texts)
    try:
        ocr = _get_ocr()
    except Exception:
        ocr = None

    blocks = []
    for i, text in enumerate(page_texts, 1):
        if len(text.strip()) >= MIN_PAGE_CHARS:
            blocks.append({"text": text, "location": f"p. {i}"})
        elif ocr is not None:
            try:
                ocr_text = ocr(path, i)
            except Exception:
                ocr_text = ""
            if ocr_text.strip():
                blocks.append({"text": ocr_text, "location": f"p. {i} (OCR)"})
            elif text.strip():
                blocks.append({"text": text, "location": f"p. {i}"})
        elif text.strip():
            blocks.append({"text": text, "location": f"p. {i}"})
    return blocks, n_pages


def _extract_binary_text(path: str) -> str:
    """Fallback text extraction for legacy binary .doc or corrupted docx files."""
    try:
        content = Path(path).read_bytes()
        ascii_strings = re.findall(rb'[\x20-\x7e\t\r\n]{4,}', content)
        text_lines = []
        for b in ascii_strings:
            try:
                s = b.decode("utf-8", errors="ignore").strip()
                if len(s) > 3 and not s.startswith("<w:") and not s.startswith("<?xml") and not s.startswith("PK"):
                    text_lines.append(s)
            except Exception:
                pass
        return "\n".join(text_lines)
    except Exception:
        return ""


def _parse_docx(path: str) -> tuple[list[dict], int]:
    blocks, buf, section = [], [], "start"

    try:
        import docx
        d = docx.Document(path)

        for para in d.paragraphs:
            t = para.text.strip()
            if not t:
                continue

            is_heading = False
            try:
                if para.style and getattr(para.style, "name", None) and str(para.style.name).startswith("Heading"):
                    is_heading = True
            except Exception:
                is_heading = False

            if is_heading:
                if buf:
                    blocks.append({"text": "\n".join(buf), "location": f"§ {section}"})
                    buf = []
                section = t[:80]
                buf.append(t)
            else:
                buf.append(t)
        if buf:
            blocks.append({"text": "\n".join(buf), "location": f"§ {section}"})
            buf = []

        for ti, table in enumerate(d.tables, 1):
            rows = []
            for row in table.rows:
                row_str = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_str:
                    rows.append(row_str)
            t = "\n".join(rows).strip()
            if t:
                blocks.append({"text": "[Table]\n" + t, "location": f"table {ti}"})
    except Exception as e:
        print(f"python-docx parsing failed for {path}: {e}")

    if not blocks:
        fallback_text = _extract_binary_text(path)
        if fallback_text.strip():
            blocks.append({"text": fallback_text, "location": "file"})

    return blocks, 0


def _parse_xlsx(path: str) -> tuple[list[dict], int]:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    blocks = []
    for ws in wb.worksheets:
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows_text.append(" | ".join(cells))
            if len(rows_text) >= 5000:
                break
        if rows_text:
            blocks.append({"text": "\n".join(rows_text), "location": f"sheet '{ws.title}'"})
    wb.close()
    return blocks, 0


def _parse_csv(path: str) -> tuple[list[dict], int]:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    return ([{"text": text, "location": "file"}] if text.strip() else []), 0


def _parse_pptx(path: str) -> tuple[list[dict], int]:
    from pptx import Presentation
    prs = Presentation(path)
    blocks = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    rows.append(" | ".join(c.text.strip() for c in row.cells))
                parts.append("[Table]\n" + "\n".join(rows))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append("[Notes] " + notes)
        if parts:
            blocks.append({"text": "\n".join(parts), "location": f"slide {i}"})
    return blocks, len(prs.slides)


def _parse_text(path: str) -> tuple[list[dict], int]:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    return ([{"text": text, "location": "file"}] if text.strip() else []), 0


def _parse_html(path: str) -> tuple[list[dict], int]:
    raw = Path(path).read_text(encoding="utf-8", errors="replace")
    # crude tag strip — good enough for indexing
    raw = re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", raw, flags=re.S | re.I)
    text = re.sub(r"<[^>]+>", " ", raw)
    text = re.sub(r"\s+", " ", text).strip()
    return ([{"text": text, "location": "file"}] if text else []), 0


def parse_file(path: str, fmt: str) -> tuple[list[dict], int]:
    """Returns (blocks, page_count)."""
    if fmt == "pdf":
        return _parse_pdf(path)
    if fmt in ("docx", "doc"):
        return _parse_docx(path)
    if fmt in ("xlsx", "xlsm"):
        return _parse_xlsx(path)
    if fmt in ("csv", "tsv"):
        return _parse_csv(path)
    if fmt == "pptx":
        return _parse_pptx(path)
    if fmt in ("html", "htm"):
        return _parse_html(path)
    return _parse_text(path)


# ---------------- chunking ----------------

def chunk_blocks(blocks: list[dict], chunk_size: int = 1500, overlap: int = 200) -> list[dict]:
    """Split blocks into ~chunk_size character chunks with overlap, keeping location labels."""
    chunk_size = max(1, int(chunk_size or 1500))
    overlap = max(0, min(int(overlap or 0), chunk_size - 1))
    chunks = []
    for block in blocks:
        text = re.sub(r"[ \t]+", " ", block["text"]).strip()
        loc = block.get("location")
        if not text:
            continue
        if len(text) <= chunk_size:
            chunks.append({"text": text, "location": loc})
            continue
        # split on paragraph boundaries where possible
        paras = re.split(r"\n{2,}|\n(?=[A-Z0-9§\[])", text)
        cur = ""
        for para in paras:
            para = para.strip()
            if not para:
                continue
            if len(para) > chunk_size:
                if cur:
                    chunks.append({"text": cur.strip(), "location": loc})
                    cur = ""
                start = 0
                while start < len(para):
                    end = min(start + chunk_size, len(para))
                    chunks.append({"text": para[start:end].strip(), "location": loc})
                    if end == len(para):
                        break
                    start = max(end - overlap, start + 1)
                continue

            candidate = f"{cur}\n\n{para}" if cur else para
            if len(candidate) <= chunk_size:
                cur = candidate
                continue

            chunks.append({"text": cur.strip(), "location": loc})
            tail = cur[-overlap:].strip() if overlap else ""
            cur = f"{tail}\n\n{para}" if tail else para
            if len(cur) > chunk_size:
                chunks.append({"text": para, "location": loc})
                cur = ""
        if cur:
            chunks.append({"text": cur.strip(), "location": loc})
    return chunks
